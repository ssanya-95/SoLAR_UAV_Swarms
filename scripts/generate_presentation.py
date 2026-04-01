from __future__ import annotations

import csv
import json
import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib-codex")

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = REPO_ROOT / "presentation"
ASSETS_DIR = PRESENTATION_DIR / "assets"
ANALYSIS_DIR = REPO_ROOT / "analysis" / "globecom_eval"
FIGURES_DIR = ANALYSIS_DIR / "figures"
TABLES_DIR = ANALYSIS_DIR / "tables"


def load_overall_summary() -> list[dict[str, str]]:
    with (TABLES_DIR / "overall_router_summary.csv").open("r", encoding="utf-8") as csv_file:
        return list(csv.DictReader(csv_file))


def load_summary_json() -> dict[str, object]:
    with (ANALYSIS_DIR / "summary.json").open("r", encoding="utf-8") as json_file:
        return json.load(json_file)


def count_significant_pdr_wins() -> int:
    count = 0
    with (TABLES_DIR / "paired_significance.csv").open("r", encoding="utf-8") as csv_file:
        for row in csv.DictReader(csv_file):
            if row["baseline"] not in {"slar-default", "gpsr", "aodv"}:
                continue
            if row["metric"] != "packet_delivery_ratio":
                continue
            if row["significant_p_lt_0_05"] != "True":
                continue
            if float(row["mean_diff"]) <= 0.0:
                continue
            count += 1
    return count


def build_overall_pdr_chart(summary_rows: list[dict[str, str]], output_path: Path) -> None:
    ordered = sorted(summary_rows, key=lambda row: float(row["packet_delivery_ratio_mean"]), reverse=True)
    routers = [row["router"] for row in ordered]
    means = [float(row["packet_delivery_ratio_mean"]) for row in ordered]
    errors = [float(row["packet_delivery_ratio_ci95"]) for row in ordered]

    colors = ["#0F766E", "#2563EB", "#7C3AED", "#6B7280", "#9CA3AF"]

    plt.figure(figsize=(8.5, 4.8))
    plt.bar(routers, means, yerr=errors, capsize=5, color=colors[: len(routers)])
    plt.ylim(0.65, 0.86)
    plt.ylabel("Packet Delivery Ratio")
    plt.title("Held-Out Overall PDR")
    for index, value in enumerate(means):
        plt.text(index, value + 0.006, f"{value:.3f}", ha="center", fontsize=10)
    plt.tight_layout()
    plt.savefig(output_path, dpi=220)
    plt.close()


def build_architecture_diagram(output_path: Path) -> None:
    fig, axis = plt.subplots(figsize=(10, 4.4))
    axis.set_xlim(0, 10)
    axis.set_ylim(0, 4.4)
    axis.axis("off")

    boxes = [
        (0.4, 2.7, 1.5, 0.9, "#DBEAFE", "Config\nParameters"),
        (2.1, 2.7, 1.7, 0.9, "#DCFCE7", "UAVNode\nMobility + Queue"),
        (4.1, 2.7, 1.7, 0.9, "#FCE7F3", "Kalman Filter\nPosition + Velocity"),
        (6.1, 2.7, 1.5, 0.9, "#FEF3C7", "Channel\nP_link"),
        (7.9, 2.7, 1.5, 0.9, "#EDE9FE", "Router\nGPSR/AODV/SLAR"),
        (3.0, 1.1, 2.0, 0.9, "#E0F2FE", "SimulationLoop\nStep, TTL, Delivery"),
        (5.4, 1.1, 2.2, 0.9, "#F3F4F6", "Analysis Scripts\nSensitivity + Eval"),
    ]

    for x_coord, y_coord, width, height, color, label in boxes:
        patch = FancyBboxPatch(
            (x_coord, y_coord),
            width,
            height,
            boxstyle="round,pad=0.02,rounding_size=0.08",
            linewidth=1.5,
            edgecolor="#334155",
            facecolor=color,
        )
        axis.add_patch(patch)
        axis.text(x_coord + width / 2, y_coord + height / 2, label, ha="center", va="center", fontsize=12)

    arrows = [
        ((1.9, 3.15), (2.1, 3.15)),
        ((3.8, 3.15), (4.1, 3.15)),
        ((5.8, 3.15), (6.1, 3.15)),
        ((7.6, 3.15), (7.9, 3.15)),
        ((4.8, 2.7), (4.2, 2.0)),
        ((7.1, 2.7), (6.4, 2.0)),
    ]
    for start, end in arrows:
        axis.annotate("", xy=end, xytext=start, arrowprops=dict(arrowstyle="->", lw=1.6, color="#334155"))

    axis.text(5.0, 4.0, "Current SLAR Framework", ha="center", va="center", fontsize=18, fontweight="bold")
    plt.tight_layout()
    plt.savefig(output_path, dpi=220, bbox_inches="tight")
    plt.close()


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 250, 252)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.75))
    text_frame = title_box.text_frame
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(15, 23, 42)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(12.0), Inches(0.4))
        sub_frame = subtitle_box.text_frame
        sub_paragraph = sub_frame.paragraphs[0]
        sub_run = sub_paragraph.add_run()
        sub_run.text = subtitle
        sub_run.font.size = Pt(14)
        sub_run.font.color.rgb = RGBColor(71, 85, 105)


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, font_size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)
        paragraph.space_after = Pt(10)


def add_section_label(slide, text: str, left: float, top: float, width: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(15, 118, 110)
    shape.line.color.rgb = RGBColor(15, 118, 110)
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_image(slide, image_path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    if height is None:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_presentation() -> Path:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    summary_rows = load_overall_summary()
    summary_json = load_summary_json()
    tuned_weights = summary_json["tuned_weights"]
    significant_wins = count_significant_pdr_wins()

    overall_pdr_path = ASSETS_DIR / "overall_pdr.png"
    architecture_path = ASSETS_DIR / "framework_architecture.png"
    build_overall_pdr_chart(summary_rows, overall_pdr_path)
    build_architecture_diagram(architecture_path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # Slide 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "SLAR UAV Swarm Routing", "Progress update for professor presentation")
    add_section_label(slide, "What this deck covers", 0.58, 1.35, 2.15)
    add_bullets(
        slide,
        [
            "Problem motivation and SLAR idea",
            "What has been implemented in the simulator",
            "Current results from sensitivity and held-out evaluation",
            "What is still missing before paper submission",
        ],
        left=0.7,
        top=1.95,
        width=6.2,
        height=3.6,
        font_size=22,
    )
    add_image(slide, overall_pdr_path, left=7.2, top=1.55, width=5.3)

    # Slide 2
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Problem and Goal")
    add_bullets(
        slide,
        [
            "UAV swarm routing is hard because links are stochastic and topology changes quickly.",
            "Classical GPSR uses only geometry; AODV uses hop count but ignores link reliability.",
            "SLAR adds two missing signals: link success probability and predicted link duration.",
            "Goal: build a complete analytical simulator and test whether this scoring improves delivery.",
        ],
        left=0.7,
        top=1.4,
        width=6.1,
        height=4.9,
        font_size=21,
    )
    add_section_label(slide, "Core scoring idea", 7.0, 1.45, 2.1)
    formula_box = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.4), Inches(1.6))
    formula_frame = formula_box.text_frame
    formula_paragraph = formula_frame.paragraphs[0]
    formula_run = formula_paragraph.add_run()
    formula_run.text = "score(j) = w_geo * geo_norm(j) + w_link * P_link(i,j) + w_ld * LD(i,j)"
    formula_run.font.size = Pt(22)
    formula_run.font.bold = True
    formula_run.font.color.rgb = RGBColor(15, 23, 42)
    add_bullets(
        slide,
        [
            "Only neighbors with strictly positive geographic progress are considered.",
            "This keeps the routing loop-free while still letting link quality change the ranking.",
        ],
        left=7.1,
        top=3.25,
        width=5.1,
        height=2.2,
        font_size=18,
    )

    # Slide 3
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "What Has Been Implemented")
    add_section_label(slide, "Simulation core", 0.65, 1.3, 1.85)
    add_bullets(
        slide,
        [
            "Stochastic channel model: log-distance path loss, temporally correlated shadowing, Rician P_link.",
            "6D Kalman filter per UAV: position + velocity, constant-velocity model, normalized link duration.",
            "Routers: GPSR, AODV, SLAR, and SLAR-geo containment baseline.",
            "Packet TTL, queueing, mobility, and delivery pipeline inside SimulationLoop.",
        ],
        left=0.75,
        top=1.9,
        width=5.9,
        height=4.8,
        font_size=20,
    )
    add_section_label(slide, "Evaluation scripts", 6.95, 1.3, 2.05)
    add_bullets(
        slide,
        [
            "Weight sensitivity analysis over the SLAR simplex.",
            "Held-out train/test evaluation across scenario families.",
            "Confidence intervals, paired significance tests, and figure/table export.",
            "Current artifacts already saved in analysis/globecom_eval/.",
        ],
        left=7.05,
        top=1.9,
        width=5.4,
        height=4.8,
        font_size=20,
    )

    # Slide 4
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Framework Architecture")
    add_image(slide, architecture_path, left=0.75, top=1.35, width=11.8)
    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.55), Inches(11.8), Inches(0.35))
    footer_frame = footer.text_frame
    footer_paragraph = footer_frame.paragraphs[0]
    footer_run = footer_paragraph.add_run()
    footer_run.text = "The workflow is now end-to-end: from analytical channel modeling to repeatable evaluation artifacts."
    footer_run.font.size = Pt(15)
    footer_run.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 5
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "SLAR Logic and Current Insight")
    add_bullets(
        slide,
        [
            "Geo term ensures forward progress and preserves GPSR-style loop avoidance.",
            "P_link captures instantaneous reliability under fading and shadowing.",
            "LD estimates whether the link will remain usable over a short horizon.",
            f"Held-out tuning currently prefers a link-heavy setting: ({tuned_weights['w_geo']:.1f}, {tuned_weights['w_link']:.1f}, {tuned_weights['w_ld']:.1f}).",
        ],
        left=0.75,
        top=1.55,
        width=6.0,
        height=4.6,
        font_size=21,
    )
    insight_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.2),
        Inches(1.8),
        Inches(5.0),
        Inches(2.2),
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(236, 253, 245)
    insight_box.line.color.rgb = RGBColor(16, 185, 129)
    frame = insight_box.text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Interpretation\n\nIn the current simulator, link reliability contributes more than LD. That suggests the main gain is coming from avoiding bad links rather than from long-horizon trajectory prediction."
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(20, 83, 45)

    # Slide 6
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Weight Sensitivity Analysis")
    add_image(slide, FIGURES_DIR / "tuning_heatmap.png", left=0.7, top=1.4, width=7.0)
    add_bullets(
        slide,
        [
            "Best training weights are in the link-heavy region.",
            "Default 1/3-1/3-1/3 is decent, but not optimal.",
            "Top held-out tuned setting: strong weight on P_link, smaller role for geo, almost no LD.",
            "This is evidence of where SLAR is getting its gains in the current model.",
        ],
        left=8.0,
        top=1.6,
        width=4.3,
        height=4.6,
        font_size=18,
    )

    # Slide 7
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Held-Out Overall Results")
    add_image(slide, overall_pdr_path, left=0.65, top=1.35, width=6.5)
    add_bullets(
        slide,
        [
            "SLAR-tuned: 0.8115 +/- 0.0297 PDR",
            "SLAR-default: 0.7878 +/- 0.0316",
            "GPSR: 0.7654 +/- 0.0313",
            "AODV: 0.7535 +/- 0.0358",
            f"Across the scenario matrix, tuned SLAR shows {significant_wins} significant PDR wins against the main baselines.",
        ],
        left=7.5,
        top=1.55,
        width=4.7,
        height=4.8,
        font_size=19,
    )

    # Slide 8
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Across Scenario Families")
    add_image(slide, FIGURES_DIR / "pdr_by_scenario.png", left=0.55, top=1.3, width=12.2)
    note_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.0),
        Inches(5.8),
        Inches(4.3),
        Inches(0.9),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
    note_box.line.color.rgb = RGBColor(59, 130, 246)
    frame = note_box.text_frame
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Takeaway: gains are not uniform, but tuned SLAR is competitive across the matrix and clearly stronger overall."
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(30, 64, 175)

    # Slide 9
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Current Status and Next Steps")
    add_section_label(slide, "What is strong now", 0.7, 1.35, 2.2)
    add_bullets(
        slide,
        [
            "End-to-end analytical simulator is implemented.",
            "Held-out evaluation pipeline exists and runs automatically.",
            "Current results show SLAR beating GPSR and AODV in the present simulator.",
        ],
        left=0.8,
        top=1.95,
        width=5.6,
        height=3.0,
        font_size=20,
    )
    add_section_label(slide, "What is still missing", 6.9, 1.35, 2.35)
    add_bullets(
        slide,
        [
            "Stronger baselines such as OLSR, P-OLSR, or predictive geographic variants.",
            "More realism and possibly a more standard simulator backend for publication-grade evidence.",
            "Final paper-quality figures, tables, and positioning for submission.",
        ],
        left=7.0,
        top=1.95,
        width=5.2,
        height=3.2,
        font_size=19,
    )
    add_section_label(slide, "Suggested ask to professor", 2.8, 5.45, 2.8)
    ask_box = slide.shapes.add_textbox(Inches(2.6), Inches(5.95), Inches(8.2), Inches(0.9))
    ask_frame = ask_box.text_frame
    ask_paragraph = ask_frame.paragraphs[0]
    ask_run = ask_paragraph.add_run()
    ask_run.text = "Request feedback on baseline selection, evaluation depth, and whether the current direction should target a stronger simulation study or a theory-first paper."
    ask_run.font.size = Pt(18)
    ask_run.font.bold = True
    ask_run.font.color.rgb = RGBColor(15, 23, 42)
    ask_paragraph.alignment = PP_ALIGN.CENTER

    output_path = PRESENTATION_DIR / "SLAR_Professor_Update.pptx"
    presentation.save(output_path)
    return output_path


def write_notes() -> Path:
    notes_path = PRESENTATION_DIR / "SLAR_Professor_Update_notes.md"
    notes = """# SLAR Professor Update Notes

## Slide 1
- This is a progress update, not a final paper defense.
- I want to show what is already implemented, what the current evidence says, and what remains.

## Slide 2
- The core motivation is that UAV links are unreliable and fast-changing.
- GPSR and AODV each miss an important part of the problem.

## Slide 3
- The main engineering milestone is that the simulator is now end-to-end.
- I am no longer discussing only theory; I can run repeatable experiments.

## Slide 4
- This slide shows how the codebase is organized.
- The important point is that the routing, channel, prediction, and evaluation layers are all separated cleanly.

## Slide 5
- The current tuned result says link reliability is the dominant factor in my simulator.
- That is useful because it tells me where the protocol is actually getting value.

## Slide 6
- The heatmap is the main sensitivity-analysis result.
- I would emphasize that the best region is link-heavy, not purely geographic.

## Slide 7
- These are held-out results, not the smoke run.
- The claim is modest: SLAR is currently better in this simulator, not yet universally proven.

## Slide 8
- The scenario-family plot shows robustness across different conditions.
- I would avoid overselling cases where the margins are small.

## Slide 9
- Be explicit that the project is in a strong prototype stage.
- Ask the professor whether to prioritize stronger baselines, simulator realism, or paper drafting next.
"""
    notes_path.write_text(notes, encoding="utf-8")
    return notes_path


def main() -> None:
    output_path = build_presentation()
    notes_path = write_notes()
    print(f"Saved presentation to {output_path}")
    print(f"Saved notes to {notes_path}")


if __name__ == "__main__":
    main()
